"""
agents/input_processor.py
─────────────────────────
Agent 1 – Input Processor

Reads every file the user supplies (sample paper, PPT, PDF, diagrams, etc.)
and returns a single structured dict that all other agents consume.

Supported formats
  • .docx  – via python-docx
  • .pdf   – via PyMuPDF (fitz)
  • .pptx  – via python-pptx
  • .txt   – plain text
  • images – stored as file path (used later by Word Generator)
"""

import os
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))


# ── helpers ──────────────────────────────────────────────────────────────────

def _read_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_pdf(path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    pages = [page.get_text() for page in doc]
    return "\n".join(pages)


def _read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _read_txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def _read_file(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    readers = {
        ".docx": _read_docx,
        ".pdf":  _read_pdf,
        ".pptx": _read_pptx,
        ".txt":  _read_txt,
    }
    if ext in readers:
        return readers[ext](path)
    # Treat as plain text as a fallback
    return _read_txt(path)


# ── public API ────────────────────────────────────────────────────────────────

def process_inputs(
    title: str,
    sample_paper_path: str,
    workflow_path: str,
    ppt_path: str,
    hardware: str,
    software: str,
    results: str,
    diagrams_description: str,
    diagram_image_paths: list[str] | None = None,
) -> dict:
    """
    Collect and normalise every user input.

    Returns
    -------
    dict with keys:
        title, sample_paper_text, workflow_text, ppt_text,
        hardware, software, results, diagrams_description,
        diagram_image_paths
    """
    print("[InputProcessor] Reading sample paper …")
    sample_text = _read_file(sample_paper_path) if sample_paper_path else ""

    print("[InputProcessor] Reading workflow / methodology …")
    workflow_text = _read_file(workflow_path) if os.path.isfile(workflow_path) else workflow_path

    print("[InputProcessor] Reading presentation …")
    ppt_text = _read_file(ppt_path) if ppt_path and os.path.isfile(ppt_path) else ""

    return {
        "title":               title,
        "sample_paper_text":   sample_text,
        "workflow_text":       workflow_text,
        "ppt_text":            ppt_text,
        "hardware":            hardware,
        "software":            software,
        "results":             results,
        "diagrams_description": diagrams_description,
        "diagram_image_paths": diagram_image_paths or [],
    }
